import pandas as pd
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from app.debug_tools import debug_wrap


@debug_wrap
def debug_wrap(func):
    def wrapper(*args, **kwargs):
        try:
            print(f"[DEBUG] Appel fonction : {func.__name__}")
            print(f"[DEBUG] Args : {args}")
            print(f"[DEBUG] Kwargs : {kwargs}")
            result = func(*args, **kwargs)
            print(f"[DEBUG] Retour fonction {func.__name__} : OK")
            return result
        except Exception as e:
            print(f"[ERROR] Exception dans {func.__name__} : {e}")
            raise
    return wrapper

@debug_wrap
def generate_ppt_from_excel(path_excel, output_ppt, llm_output=None):
    df = pd.read_excel(path_excel)

    required_cols = ["semaines", "Prévu cumulé (€)", "Réalisé cumulé (€)", "Écart cumulé (€)"]
    if not all(col in df.columns for col in required_cols):
        raise ValueError(f"Colonnes manquantes dans l'Excel. Colonnes attendues : {required_cols}")

    df = df.replace("—", pd.NA)

    def clean_number(x):
        if pd.isna(x):
            return pd.NA
        x = str(x)
        x = x.replace("€", "").replace(" ", "").replace("\u202f", "")
        x = x.replace(",", ".").replace("—", "")
        return x

    numeric_cols = ["Prévu cumulé (€)", "Réalisé cumulé (€)", "Écart cumulé (€)"]
    for col in numeric_cols:
        df[col] = df[col].apply(clean_number)
        df[col] = pd.to_numeric(df[col], errors="coerce")

    semaines = df["semaines"].astype(str)
    prevu = df["Prévu cumulé (€)"]
    realise = df["Réalisé cumulé (€)"]
    ecart = df["Écart cumulé (€)"]

    graph_path = "graph_temp.png"
    plt.figure(figsize=(10, 5))
    plt.plot(semaines, prevu, label="Prévu cumulé", color="blue", marker="o")
    plt.plot(semaines, realise, label="Réalisé cumulé", color="green", marker="o")
    plt.plot(semaines, ecart, label="Écart cumulé", color="red", linestyle="--", marker="x")
    plt.legend()
    plt.grid(True)
    plt.tight_layout()
    plt.savefig(graph_path)
    plt.close()

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title_shape = slide.shapes.title
    title_shape.text = llm_output["title"] if llm_output else "Suivi Budget – Activité Développement"

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(1))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = llm_output["subtitle"] if llm_output else "Prévu vs Réalisé – Semaine en cours"
    subtitle_tf.paragraphs[0].font.size = Pt(18)

    slide.shapes.add_picture(graph_path, Inches(1), Inches(2), width=Inches(8))

    txBox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.text = llm_output["key_message"] if llm_output else "Principales justifications :"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(16)

    justifications = llm_output["justifications"] if llm_output else [
        "Sous-staffing ponctuel",
        "Démarrage opérationnel lent",
        "Dépendances externes"
    ]

    for j in justifications:
        p = tf.add_paragraph()
        p.text = j
        p.level = 1
        p.font.size = Pt(14)

    prs.save(output_ppt)

    if os.path.exists(graph_path):
        os.remove(graph_path)
